import os
import json
import urllib.request
import urllib.error
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv
from llama_cloud_services import LlamaParse
from prompt_templates import ENRICH_PRESENTATION_PROMPT

# Load environment variables from .env file
load_dotenv()

LLAMA_CLOUD_API_KEY = os.getenv("LLAMA_CLOUD_API_KEY")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")

if not LLAMA_CLOUD_API_KEY:
    raise ValueError("LLAMA_CLOUD_API_KEY must be set in the environment.")
if not ANTHROPIC_API_KEY:
    raise ValueError("ANTHROPIC_API_KEY must be set in the environment.")

base_dir = os.path.dirname(os.path.abspath(__file__))
intermediate_dir = os.path.join(base_dir, "intermediate")
os.makedirs(intermediate_dir, exist_ok=True)

def make_api_call(api_key, content):
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        'x-api-key': api_key,
        'anthropic-version': '2023-06-01',
        'content-type': 'application/json'
    }
    data = json.dumps({
        "model": "claude-3-5-sonnet-20241022",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": content}]
    }).encode('utf-8')

    req = urllib.request.Request(url, data=data, headers=headers, method='POST')

    try:
        with urllib.request.urlopen(req) as response:
            output = response.read()
            return json.loads(output)['content'][0]['text']
    except Exception as e:
        print(f"API error: {e}")
    return None

def extract_document_data(doc_path):
    parser = LlamaParse(api_key=LLAMA_CLOUD_API_KEY, language="en", verbose=True)
    result = parser.parse(doc_path)

    text_documents = result.get_text_documents(split_by_page=False)
    text = "\n".join(doc.text for doc in text_documents)

    images_dir = os.path.join(base_dir, "images")
    if os.path.exists(images_dir):
        for f in os.listdir(images_dir):
            file_path = os.path.join(images_dir, f)
            if os.path.isfile(file_path):
                os.remove(file_path)
    else:
        os.makedirs(images_dir, exist_ok=True)

    try:
        images = result.get_image_documents(
            include_screenshot_images=True,
            include_object_images=True,
            image_download_dir=images_dir
        )
    except Exception as e:
        print(f"[WARNING] Failed to fetch images from LlamaParse: {e}")
        images = []

    image_paths = []
    for idx, image_doc in enumerate(images):
        image_path = image_doc.image_path
        page_index = getattr(image_doc, 'page_index', idx)
        context_text = getattr(image_doc, 'context_text', '')

        if not context_text and page_index < len(result.pages):
            context_text = result.pages[page_index].text.lower().strip()[:100]

        context_text = context_text.replace(" ", "_").replace("/", "_")
        safe_context = ''.join(c for c in context_text if c.isalnum() or c == '_')
        new_filename = f"img_p{page_index}_{idx}_{safe_context}.png"
        new_path = os.path.join(images_dir, new_filename)

        try:
            os.rename(image_path, new_path)
            image_paths.append(new_path)
            print(f"[DEBUG] Renamed image to: {new_filename}")
        except Exception as e:
            print(f"[WARNING] Failed to rename image {image_path}: {e}")

    print(f"[DEBUG] Number of images extracted: {len(image_paths)}")

    tables = []
    for page in result.pages:
        if hasattr(page, "structuredData") and page.structuredData:
            if "tables" in page.structuredData:
                tables.extend(page.structuredData["tables"])

    print(f"[DEBUG] Number of tables extracted: {len(tables)}")

    return {"text": text, "images": image_paths, "tables": tables}

def enrich_with_claude(document_data):
    print(f"[DEBUG] Table data sent to Claude:\n{json.dumps(document_data['tables'], indent=2)}")
    print(f"[DEBUG] Image data sent to Claude:\n{json.dumps(document_data['images'], indent=2)}")

    formatted_prompt = ENRICH_PRESENTATION_PROMPT.format(
        text_content=document_data['text'],
        image_paths=document_data['images'],
        table_data=json.dumps(document_data['tables'], indent=2)
    )

    response = make_api_call(ANTHROPIC_API_KEY, formatted_prompt)

    if response is None:
        raise ValueError("Claude API call failed. No response received.")

    try:
        structured_response = json.loads(response)
    except json.JSONDecodeError:
        raise ValueError("Claude response is not valid JSON. Response was:\n" + response)

    return structured_response

def add_table_to_slide(slide, table_data):
    if not table_data or not isinstance(table_data, list) or not table_data[0]:
        print("No valid table data provided.")
        return

    rows, cols = len(table_data), len(table_data[0])

    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(7.0)
    height = Inches(2.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    total_width = Inches(7.0)
    col_width = int(total_width / cols)
    for col in range(cols):
        table.columns[col].width = col_width

    for i, row in enumerate(table_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True if i == 0 else False
            cell.text_frame.paragraphs[0].alignment = 1

    print("Table added successfully!")

def create_ppt_from_claude(claude_slides, document_data, ppt_path="generated_deck.pptx", template_path=None):
    ppt = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()
    if template_path and len(ppt.slides) > 0:
        rId = ppt.slides._sldIdLst[0].rId
        ppt.part.drop_rel(rId)
        del ppt.slides._sldIdLst[0]

    for slide_data in claude_slides:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        if slide_data.get("title") and slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]
        if slide_data.get("text") and len(slide.placeholders) > 1:
            text_box = slide.placeholders[1]
            text_frame = text_box.text_frame
            text_frame.clear()
            for idx, line in enumerate(slide_data["text"].split("\n")):
                if not line.strip(): continue
                if idx == 0:
                    text_frame.text = line
                    text_frame.paragraphs[0].bullet = True
                else:
                    p = text_frame.add_paragraph()
                    p.text = line
                    p.bullet = True

        if slide_data.get("table"):
            print(f"[DEBUG] Slide titled '{slide_data.get('title')}' contains a table with {len(slide_data['table'])} rows")
            add_table_to_slide(slide, slide_data["table"])
        else:
            print(f"[DEBUG] Slide titled '{slide_data.get('title')}' has NO table")

        assigned_image = None
        slide_title = slide_data.get("title", "").lower()

        for image_path in document_data["images"]:
            image_filename = os.path.basename(image_path).lower()
            if any(keyword in image_filename for keyword in slide_title.split()):
                assigned_image = image_path
                document_data["images"].remove(image_path)
                break

        if assigned_image:
            slide_data["image"] = assigned_image
        else:
            print("No contextual match found, skipping image assignment for this slide.")

        if slide_data.get("image") and os.path.exists(slide_data["image"]):
            img_path = slide_data["image"]
            slide_width = ppt.slide_width
            slide_height = ppt.slide_height
            margin = Inches(0.3)
            image_width = Inches(4.5)
            image_left = slide_width - image_width - margin
            image_height = Inches(3)
            image_top = (slide_height - image_height) / 2
            slide.shapes.add_picture(img_path, image_left, image_top, width=image_width, height=image_height)

    thank_you_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    if thank_you_slide.shapes.title:
        thank_you_slide.shapes.title.text = "Thank You!"

    if len(thank_you_slide.placeholders) > 1:
        text_box = thank_you_slide.placeholders[1]
        text_frame = text_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = "We appreciate your attention. Looking forward to your questions!"

    ppt.save(ppt_path)
    print(f"Presentation saved as {ppt_path}")

def main():
    doc_path = os.path.join(base_dir, "input", "doc.docx")
    output_ppt_path = os.path.join(base_dir, "output", "output_presentation.pptx")
    template_path = os.path.join(base_dir, "template", "template.pptx")

    document_data = extract_document_data(doc_path)

    extracted_data_path = os.path.join(intermediate_dir, "extracted_data.json")
    with open(extracted_data_path, "w", encoding="utf-8") as json_file:
        json.dump(document_data, json_file, ensure_ascii=False, indent=4)

    claude_response = enrich_with_claude(document_data)

    claude_response_path = os.path.join(intermediate_dir, "claude_structured_response.json")
    with open(claude_response_path, "w", encoding="utf-8") as json_file:
        json.dump(claude_response, json_file, ensure_ascii=False, indent=4)

    create_ppt_from_claude(claude_response, document_data, output_ppt_path, template_path)

if __name__ == "__main__":
    main()

